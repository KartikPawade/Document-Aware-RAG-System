"""
ingestion/parsers/office_plugins.py
─────────────────────────────────────
Plugins for: PPTX, XLSX, CSV, TXT, Markdown.
Each implements BaseFormatPlugin independently.
"""
from __future__ import annotations

import io
from pathlib import Path
from typing import List, Optional

import pandas as pd

from src.core.logging import get_logger
from src.core.models import (
    Chunk, ContentType, DocumentFormat,
    ImageBlock, ParsedDocument, TextBlock,
)
from src.ingestion.chunkers.utils import (
    create_parent_child_chunks,
    heading_based_chunk,
    semantic_chunk,
    estimate_tokens,
)
from src.ingestion.parsers.base import BaseFormatPlugin, ChunkConfig

logger = get_logger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# PPTX Plugin
# ─────────────────────────────────────────────────────────────────────────────

class PPTXPlugin(BaseFormatPlugin):
    """
    PowerPoint parser: one chunk per slide + speaker notes.
    Slide title becomes the section heading.
    """
    supported_extensions = [".pptx", ".ppt"]
    document_format = DocumentFormat.PPTX

    def parse(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        from pptx import Presentation

        text_blocks: List[TextBlock] = []
        structure = {"slide_titles": [], "slide_count": 0}

        try:
            prs = Presentation(io.BytesIO(file_bytes))
            structure["slide_count"] = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, start=1):
                # Extract title
                title = ""
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
                    structure["slide_titles"].append(
                        {"slide": slide_num, "title": title}
                    )

                # Extract all text from shapes
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = " ".join(
                                run.text for run in para.runs
                            ).strip()
                            if para_text and para_text != title:
                                slide_texts.append(para_text)

                # Speaker notes
                notes_text = ""
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_text = notes_tf.text.strip() if notes_tf else ""

                slide_content = "\n".join(slide_texts)
                if notes_text:
                    slide_content += f"\n\nSpeaker Notes: {notes_text}"

                if slide_content.strip():
                    text_blocks.append(TextBlock(
                        text=slide_content,
                        heading=title or f"Slide {slide_num}",
                        slide_number=slide_num,
                    ))

        except Exception as e:
            logger.error("PPTX parsing failed", filename=filename, error=str(e))

        return ParsedDocument(
            text_blocks=text_blocks,
            tables=[],
            images=[],
            structure=structure,
            source_metadata={"filename": filename, "format": "pptx"},
        )

    def chunk(self, doc: ParsedDocument, config: ChunkConfig) -> List[Chunk]:
        """One chunk per slide — slide is a natural semantic unit."""
        chunks: List[Chunk] = []
        for block in doc.text_blocks:
            chunk = Chunk(
                text=block.text,
                section_heading=block.heading,
                slide_number=block.slide_number,
                source_id=doc.source_metadata.get("source_id", ""),
                format=DocumentFormat.PPTX,
                content_type=ContentType.SLIDE,
                token_count=estimate_tokens(block.text),
            )
            chunks.append(chunk)
        return chunks


# ─────────────────────────────────────────────────────────────────────────────
# XLSX Plugin
# ─────────────────────────────────────────────────────────────────────────────

class XLSXPlugin(BaseFormatPlugin):
    """
    Excel parser.
    IMPORTANT: Rows are NOT embedded. They are stored in PostgreSQL.
    Only table schema + sample rows are chunked for routing.
    """
    supported_extensions = [".xlsx", ".xls", ".xlsm"]
    document_format = DocumentFormat.XLSX

    def parse(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        tables: List[pd.DataFrame] = []
        text_blocks: List[TextBlock] = []

        try:
            xl = pd.ExcelFile(io.BytesIO(file_bytes))
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                df = df.dropna(how="all")
                if not df.empty:
                    df.attrs["sheet_name"] = sheet_name
                    df.attrs["source_filename"] = filename
                    tables.append(df)

        except Exception as e:
            logger.error("XLSX parsing failed", filename=filename, error=str(e))

        return ParsedDocument(
            text_blocks=text_blocks,
            tables=tables,
            images=[],
            structure={"sheet_count": len(tables)},
            source_metadata={"filename": filename, "format": "xlsx"},
        )

    def chunk(self, doc: ParsedDocument, config: ChunkConfig) -> List[Chunk]:
        """Generate schema chunks only — actual data goes to SQL store."""
        chunks: List[Chunk] = []
        for df in doc.tables:
            sheet_name = df.attrs.get("sheet_name", "Sheet1")
            filename = doc.source_metadata.get("filename", "")
            schema_text = self._build_schema_chunk(df, sheet_name, filename)
            chunk = Chunk(
                text=schema_text,
                section_heading=f"Table: {sheet_name}",
                source_id=doc.source_metadata.get("source_id", ""),
                format=DocumentFormat.XLSX,
                content_type=ContentType.TABLE_REF,
                token_count=estimate_tokens(schema_text),
            )
            chunks.append(chunk)
        return chunks

    def _build_schema_chunk(self, df: pd.DataFrame, sheet_name: str, filename: str) -> str:
        cols = ", ".join(f"{c} ({df[c].dtype})" for c in df.columns)
        sample = df.head(3).to_string(index=False)
        row_count = len(df)
        return (
            f"Spreadsheet: {filename}, Sheet: {sheet_name}. "
            f"Rows: {row_count}. Columns: {cols}.\n"
            f"Sample data:\n{sample}"
        )


# ─────────────────────────────────────────────────────────────────────────────
# CSV Plugin
# ─────────────────────────────────────────────────────────────────────────────

class CSVPlugin(BaseFormatPlugin):
    """
    CSV parser — same philosophy as XLSX.
    Data goes to SQL store. Schema chunk goes to vector store.
    """
    supported_extensions = [".csv", ".tsv"]
    document_format = DocumentFormat.CSV

    def parse(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        tables: List[pd.DataFrame] = []
        sep = "\t" if filename.endswith(".tsv") else ","

        try:
            df = pd.read_csv(io.BytesIO(file_bytes), sep=sep, on_bad_lines="skip")
            df = df.dropna(how="all")
            if not df.empty:
                df.attrs["source_filename"] = filename
                tables.append(df)
        except Exception as e:
            logger.error("CSV parsing failed", filename=filename, error=str(e))

        return ParsedDocument(
            text_blocks=[],
            tables=tables,
            images=[],
            structure={"row_count": len(tables[0]) if tables else 0},
            source_metadata={"filename": filename, "format": "csv"},
        )

    def chunk(self, doc: ParsedDocument, config: ChunkConfig) -> List[Chunk]:
        chunks: List[Chunk] = []
        for df in doc.tables:
            filename = doc.source_metadata.get("filename", "")
            cols = ", ".join(f"{c} ({df[c].dtype})" for c in df.columns)
            sample = df.head(3).to_string(index=False)
            schema_text = (
                f"CSV file: {filename}. "
                f"Rows: {len(df)}. Columns: {cols}.\n"
                f"Sample:\n{sample}"
            )
            chunk = Chunk(
                text=schema_text,
                source_id=doc.source_metadata.get("source_id", ""),
                format=DocumentFormat.CSV,
                content_type=ContentType.TABLE_REF,
                token_count=estimate_tokens(schema_text),
            )
            chunks.append(chunk)
        return chunks


# ─────────────────────────────────────────────────────────────────────────────
# TXT Plugin
# ─────────────────────────────────────────────────────────────────────────────

class TXTPlugin(BaseFormatPlugin):
    """Plain text parser with semantic chunking."""
    supported_extensions = [".txt", ".text", ".log"]
    document_format = DocumentFormat.TXT

    def parse(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        try:
            text = file_bytes.decode("utf-8", errors="replace")
        except Exception:
            text = ""

        blocks = [
            TextBlock(text=para.strip())
            for para in text.split("\n\n")
            if para.strip()
        ]

        return ParsedDocument(
            text_blocks=blocks,
            tables=[],
            images=[],
            structure={"paragraph_count": len(blocks)},
            source_metadata={"filename": filename, "format": "txt"},
        )

    def chunk(self, doc: ParsedDocument, config: ChunkConfig) -> List[Chunk]:
        full_text = "\n\n".join(b.text for b in doc.text_blocks)
        child_texts = semantic_chunk(full_text, config.child_chunk_size, config.chunk_overlap)

        chunks: List[Chunk] = []
        for i, text in enumerate(child_texts):
            chunk = Chunk(
                text=text,
                source_id=doc.source_metadata.get("source_id", ""),
                format=DocumentFormat.TXT,
                content_type=ContentType.PROSE,
                chunk_index=i,
                total_chunks=len(child_texts),
                token_count=estimate_tokens(text),
            )
            chunks.append(chunk)
        return chunks


# ─────────────────────────────────────────────────────────────────────────────
# Markdown Plugin
# ─────────────────────────────────────────────────────────────────────────────

class MarkdownPlugin(BaseFormatPlugin):
    """
    Markdown parser with heading-based hierarchical chunking.
    H1/H2/H3 headings define natural chunk boundaries.
    Code blocks are preserved as intact chunks.
    """
    supported_extensions = [".md", ".markdown", ".mdx"]
    document_format = DocumentFormat.MARKDOWN

    def parse(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        try:
            text = file_bytes.decode("utf-8", errors="replace")
        except Exception:
            text = ""

        text_blocks: List[TextBlock] = []
        current_heading: Optional[str] = None
        current_lines: List[str] = []
        in_code_block = False

        lines = text.split("\n")
        for line in lines:
            # Toggle code block
            if line.startswith("```"):
                in_code_block = not in_code_block
                current_lines.append(line)
                if not in_code_block and current_lines:
                    # Flush code block as a single chunk
                    code_text = "\n".join(current_lines)
                    text_blocks.append(TextBlock(
                        text=code_text,
                        heading=current_heading,
                        is_code=True,
                    ))
                    current_lines = []
                continue

            if in_code_block:
                current_lines.append(line)
                continue

            # Detect headings
            heading_match = None
            for level in (1, 2, 3):
                if line.startswith("#" * level + " ") and not line.startswith("#" * (level + 1)):
                    heading_match = line.lstrip("#").strip()
                    break

            if heading_match:
                # Flush previous content
                if current_lines:
                    content = "\n".join(current_lines).strip()
                    if content:
                        text_blocks.append(TextBlock(
                            text=content,
                            heading=current_heading,
                        ))
                current_heading = heading_match
                current_lines = []
            else:
                current_lines.append(line)

        # Flush remaining
        if current_lines:
            content = "\n".join(current_lines).strip()
            if content:
                text_blocks.append(TextBlock(text=content, heading=current_heading))

        return ParsedDocument(
            text_blocks=text_blocks,
            tables=[],
            images=[],
            structure={"headings": [b.heading for b in text_blocks if b.heading]},
            source_metadata={"filename": filename, "format": "markdown"},
        )

    def chunk(self, doc: ParsedDocument, config: ChunkConfig) -> List[Chunk]:
        chunks: List[Chunk] = []
        chunk_index = 0

        for block in doc.text_blocks:
            if block.is_code:
                # Code blocks: never split, preserve as single chunk
                chunk = Chunk(
                    text=block.text,
                    section_heading=block.heading,
                    source_id=doc.source_metadata.get("source_id", ""),
                    format=DocumentFormat.MARKDOWN,
                    content_type=ContentType.CODE,
                    chunk_index=chunk_index,
                    token_count=estimate_tokens(block.text),
                )
                chunks.append(chunk)
                chunk_index += 1
            else:
                parent, children = create_parent_child_chunks(
                    parent_text=block.text,
                    heading=block.heading,
                    source_metadata=doc.source_metadata,
                    child_target_tokens=config.child_chunk_size,
                    child_overlap_tokens=config.chunk_overlap,
                    chunk_index_offset=chunk_index,
                )
                parent.format = DocumentFormat.MARKDOWN
                for child in children:
                    child.format = DocumentFormat.MARKDOWN
                    child.content_type = ContentType.PROSE

                chunks.append(parent)
                chunks.extend(children)
                chunk_index += len(children)

        return chunks